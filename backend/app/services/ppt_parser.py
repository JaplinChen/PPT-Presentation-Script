from pptx import Presentation
from typing import List, Dict
from pathlib import Path
import time

class PPTParser:
    """解析 PowerPoint 文件並提取結構化內容 - 極限優化版"""
    
    def parse(self, ppt_path: str) -> List[Dict]:
        """
        解析 PPT 並提取所有投影片的內容
        """
        start_time = time.time()
        
        if not Path(ppt_path).exists():
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")
        
        try:
            # 讀取 Presentation 是最耗時的一步 (I/O 密集)
            prs = Presentation(ppt_path)
        except Exception as e:
            print(f"[PPTParser] ERROR: 無法讀取 PPT 檔案 - {e}")
            raise
            
        slides_data = []
        slide_height = prs.slide_height
        visible_slide_no = 0
        
        # 預計算 threshold
        threshold = slide_height * 0.25
        
        for i, slide in enumerate(prs.slides):
            # 檢查是否為隱藏投影片
            if slide.element.get('show') == '0' or slide.element.get('show') == 'false':
                continue
            
            visible_slide_no += 1
            
            title = ""
            bullets = []
            tables = []
            image_count = 0
            candidate_titles = []
            
            # 單次遍歷所有物件，減少屬性訪問次數
            for shape in slide.shapes:
                try:
                    stype = shape.shape_type
                except: continue
                
                # 1. 圖片計數 (快捷路徑)
                if stype == 13: # Picture
                    image_count += 1
                    continue
                
                # 2. 表格提取 (19 = Table)
                if stype == 19 or shape.has_table:
                    try:
                        table = shape.table
                        table_data = {
                            'rows': len(table.rows),
                            'cols': len(table.columns),
                            'content': [[cell.text.strip() for cell in row.cells] for row in table.rows]
                        }
                        tables.append(table_data)
                    except: pass
                    continue
                
                # 3. 文字內容提取
                # 僅針對有文字框的物件 (Placeholder=14, Textbox=17, AutoShape=1)
                if shape.has_text_frame:
                    try:
                        text = shape.text.strip()
                        if not text:
                            continue
                            
                        # 檢查標題佔位符 (Placeholder type 1 = Title)
                        is_title = False
                        if shape.is_placeholder:
                            ph_type = shape.placeholder_format.type
                            if ph_type == 1 or ph_type == 3: # 1=Title, 3=Centered Title
                                if not title: title = text
                                is_title = True
                        
                        if not is_title:
                            # 收集候選標題 (位於上方且文字較大)
                            if hasattr(shape, "top") and shape.top < threshold and len(text) > 1:
                                fsize = 0
                                try:
                                    if shape.text_frame.paragraphs[0].runs:
                                        fsize = shape.text_frame.paragraphs[0].runs[0].font.size or 0
                                except: pass
                                candidate_titles.append({"text": text, "top": shape.top, "fsize": fsize})
                            
                            # 收集為內容要點 (過濾極短文字如頁碼)
                            if len(text) > 1:
                                lines = [l.strip() for l in text.split('\n') if l.strip()]
                                bullets.extend(lines)
                    except: pass
            
            # 挑選最佳標題
            if not title and candidate_titles:
                # 依字體大小 > 位置排序
                candidate_titles.sort(key=lambda x: (-x['fsize'], x['top']))
                title = candidate_titles[0]['text']
            
            # 備註提取 (僅在需要時訪問，這一步有時很慢)
            notes = ""
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
            except: pass
            
            slides_data.append({
                'slide_no': visible_slide_no,
                'title': title,
                'bullets': bullets,
                'tables': tables,
                'notes': notes,
                'image_count': image_count
            })
            
            # 每 50 頁列印一次進度，避免大檔案讓用戶覺得死機
            if visible_slide_no % 50 == 0:
                print(f"[PPTParser] 已處理 {visible_slide_no} 頁...")
        
        elapsed = time.time() - start_time
        print(f"[PPTParser] 完成! 總計 {len(slides_data)} 頁, 耗時 {elapsed:.2f}s")
        return slides_data
    
    def get_summary(self, slides_data: List[Dict]) -> Dict:
        """取得 PPT 摘要資訊 (O(N) 遍歷)"""
        return {
            'total_slides': len(slides_data),
            'titled_slides': sum(1 for s in slides_data if s['title']),
            'total_bullets': sum(len(s['bullets']) for s in slides_data),
            'slides_with_notes': sum(1 for s in slides_data if s['notes']),
            'slides_with_tables': sum(1 for s in slides_data if s['tables']),
            'total_images': sum(s['image_count'] for s in slides_data)
        }
