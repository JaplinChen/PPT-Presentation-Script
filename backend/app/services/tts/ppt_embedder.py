"""
PPT audio embedding functionality.
"""
import os
import re
import shutil
import uuid
from pathlib import Path
from typing import Dict, List, Optional

from mutagen.mp3 import MP3
from pptx import Presentation
from pptx.util import Cm

class PPTEmbedder:
    """Handles embedding audio into PowerPoint presentations"""
    
    def __init__(self, output_dir: Path):
        self.output_dir = output_dir
    
    async def embed_audio(
        self,
        original_pptx_path: str,
        slide_scripts: List[Dict],
        audio_generator,  # AudioGenerator instance
        voice: str,
        rate: str = "+0%",
        pitch: str = "+0Hz",
        progress_callback: Optional[callable] = None,
    ) -> tuple[str, Dict[int, str]]:
        """
        Embed audio into PPT slides.
        
        Args:
            original_pptx_path: Path to original PPT
            slide_scripts: List of slide script data
            audio_generator: AudioGenerator instance for generating audio
            voice: TTS voice to use
            rate: Speech rate
            pitch: Speech pitch
            progress_callback: Progress reporting function
            
        Returns:
            Tuple of (output_path, slide_scripts_dict)
        """
        # Create copy
        original_path = Path(original_pptx_path)
        output_filename = f"narrated_{uuid.uuid4()}.pptx"
        output_path = self.output_dir / output_filename
        shutil.copy(original_path, output_path)
        
        prs = Presentation(output_path)
        
        # Map slide data
        script_data_map = {int(item['slide_no']): item for item in slide_scripts}
        all_slide_scripts: Dict[int, str] = {}
        
        # Process each slide
        visible_slide_index = 0
        total_slides = len(prs.slides)
        
        for i, slide in enumerate(prs.slides):
            # Skip hidden slides
            if slide.element.get('show') == '0' or slide.element.get('show') == 'false':
                continue
            
            visible_slide_index += 1
            
            if progress_callback:
                progress_callback(
                    int((visible_slide_index / total_slides) * 80), 
                    f"Processing slide {visible_slide_index}/{total_slides}..."
                )

            slide_data = script_data_map.get(visible_slide_index)
            
            
            if not slide_data:
                print(f"[PPTEmbedder] Slide {i+1} - No script data")
                continue
            
            script_text = slide_data.get('script', '')
            all_slide_scripts[i + 1] = script_text
            
            if not script_text:
                print(f"[PPTEmbedder] Slide {i+1} - No script")
                continue
            
            try:
                # Clean text: remove formatting markers and special characters
                clean_text = script_text
                
                # Remove section markers like "=== Opening ===" or "--- Slide X ---"
                clean_text = re.sub(r'===.*?===', '', clean_text)
                clean_text = re.sub(r'---.*?---', '', clean_text)
                
                # Remove time markers like "(約 8 秒)" or "(30秒)"
                clean_text = re.sub(r'\([約大概]*\s*\d+\s*[秒分鐘seconds]+\)', '', clean_text)
                
                # Remove special characters that might cause TTS issues
                clean_text = re.sub(r'[*()[\]/]', ' ', clean_text)
                
                # Remove extra whitespace
                clean_text = re.sub(r'\s+', ' ', clean_text).strip()
                
                if not clean_text.strip():
                    continue
                
                audio_info = await audio_generator.generate_audio(clean_text, voice, rate, pitch)
                audio_path = audio_info['path']
                
                # Get duration
                audio = MP3(audio_path)
                duration_sec = audio.info.length
                
                # Embed audio
                self._add_audio_to_slide(slide, audio_path, prs.slide_height, duration_sec)
                
                print(f"[PPTEmbedder] Slide {i+1} - Audio embedded")
            
            except Exception as e:
                print(f"[PPTEmbedder] Slide {i+1} - Failed: {e}")
                import traceback
                traceback.print_exc()
        
        # Save
        prs.save(output_path)
        print(f"[PPTEmbedder] Saved to {output_path}")
        
        return str(output_path.resolve()), all_slide_scripts
    
    def _add_audio_to_slide(self, slide, audio_path: str, slide_height, duration_sec: float):
        """Add audio shape to slide with auto-play"""
        # Position & Size
        icon_width = Cm(0.5)
        icon_height = Cm(0.5)
        left_pos = Cm(0.5)
        top_pos = slide_height - icon_height - Cm(0.5)
        
        # Custom icon
        backend_dir = Path("d:/Works/PPT_Presentation_Script/backend")
        asset_path = backend_dir / "assets" / "audio_icon.png"
        poster_frame = str(asset_path) if asset_path.exists() else None

        # Insert audio
        movie = slide.shapes.add_movie(
            audio_path, 
            left=left_pos, top=top_pos, 
            width=icon_width, height=icon_height, 
            poster_frame_image=poster_frame, 
            mime_type='audio/mp3'
        )
        
        # Set auto-play
        try:
            shape_id = movie.shape_id
            self._add_autoplay_timing(slide, shape_id)
        except Exception as e:
            print(f"[PPTEmbedder] Warning: Could not set auto-play - {e}")
        
        # Set transition
        if hasattr(slide, 'slide_show_transition'):
            trans = slide.slide_show_transition
            trans.advance_on_time = True
            trans.advance_after_time = int((duration_sec + 2.0) * 1000)
    
    def _add_autoplay_timing(self, slide, shape_id):
        """Inject XML for auto-play"""
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import qn
        
        timing = slide.element.find(qn('p:timing'))
        if timing is None:
            timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            extLst = slide.element.find(qn('p:extLst'))
            if extLst is not None:
                slide.element.insert(slide.element.index(extLst), timing)
            else:
                slide.element.append(timing)
        
        tnLst = timing.find(qn('p:tnLst'))
        if tnLst is None:
            tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            timing.append(tnLst)
        
        for child in list(tnLst):
            tnLst.remove(child)
        
        xml = f"""
        <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:cmd type="call" cmd="playFrom(0.0)">
                            <p:cBhvr>
                              <p:cTn id="4" dur="1" fill="hold"/>
                              <p:tgtEl>
                                <p:spTgt spid="{shape_id}"/>
                              </p:tgtEl>
                            </p:cBhvr>
                          </p:cmd>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
        """
        new_node = parse_xml(xml)
        tnLst.append(new_node)
